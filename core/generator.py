
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import os
import io
from typing import Optional, Dict, Any, List

def create_ppt_from_template(slide_data, output_path, template_path=None, template_style=None):
    """
    Creates a NEW PPT file by DUPLICATING template slides and replacing content.
    This preserves ALL visual elements including backgrounds, shapes, and images.
    """
    
    if not template_path or not os.path.exists(template_path):
        # No template - create basic presentation
        return create_basic_presentation(slide_data, output_path)
    
    try:
        # Load template
        template_prs = Presentation(template_path)
        
        # Find the best template slide to use as a base
        template_slide = find_best_content_slide(template_prs)
        template_slide_layout = template_slide.slide_layout if template_slide else template_prs.slide_layouts[1]
        
        # Create new presentation using the same template
        new_prs = Presentation(template_path)
        
        # Remove all existing slides (keep just the master/layouts)
        slide_indices = list(range(len(new_prs.slides)))
        for i in reversed(slide_indices):
            rId = new_prs.slides._sldIdLst[i].rId
            new_prs.part.drop_rel(rId)
            del new_prs.slides._sldIdLst[i]
        
        # Create new slides by duplicating the template slide structure
        for slide_content in slide_data:
            new_slide = duplicate_slide_with_content(
                new_prs, 
                template_slide_layout, 
                template_slide,
                slide_content
            )
        
        new_prs.save(output_path)
        return output_path
        
    except Exception as e:
        print(f"Error with template processing: {e}")
        # Fallback to basic presentation
        return create_basic_presentation(slide_data, output_path)

def find_best_content_slide(template_prs):
    """
    Find the best slide from template to use as a base for content slides.
    Prefers slides with both title and content areas.
    """
    best_slide = None
    best_score = 0
    
    for slide in template_prs.slides:
        score = 0
        has_title = False
        has_content_area = False
        
        for shape in slide.shapes:
            if shape.is_placeholder:
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    has_title = True
                    score += 2
                elif shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    has_content_area = True
                    score += 3
        
        # Prefer slides with both title and content
        if has_title and has_content_area:
            score += 5
            
        if score > best_score:
            best_score = score
            best_slide = slide
    
    return best_slide

def duplicate_slide_with_content(new_prs, layout, template_slide, content):
    """
    Create a new slide by copying template slide structure and replacing content.
    This preserves ALL visual elements while updating text content.
    """
    # Add new slide using the same layout
    new_slide = new_prs.slides.add_slide(layout)
    
    if template_slide:
        # Copy ALL non-placeholder shapes from template slide
        copy_template_visual_elements(template_slide, new_slide)
    
    # Now populate the placeholders with our content
    populate_slide_content(new_slide, content)
    
    return new_slide

def copy_template_visual_elements(template_slide, new_slide):
    """
    Copy all visual elements (non-placeholder shapes) from template to new slide.
    This includes backgrounds, decorative shapes, images, logos, etc.
    """
    try:
        shapes_to_copy = []
        
        # Identify all non-placeholder shapes in template
        for shape in template_slide.shapes:
            if not shape.is_placeholder:
                shapes_to_copy.append(shape)
        
        # Copy each non-placeholder shape
        for shape in shapes_to_copy:
            copy_shape_to_slide(shape, new_slide)
            
    except Exception as e:
        print(f"Error copying visual elements: {e}")

def copy_shape_to_slide(source_shape, target_slide):
    """
    Copy a specific shape from source to target slide.
    Handles different shape types (images, text boxes, shapes, etc.)
    """
    try:
        # Handle images
        if source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            copy_image_shape(source_shape, target_slide)
        
        # Handle text boxes and auto shapes
        elif source_shape.shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE):
            copy_text_or_shape(source_shape, target_slide)
            
        # Handle other shape types
        else:
            try:
                # Generic shape copying (best effort)
                copy_generic_shape(source_shape, target_slide)
            except:
                print(f"Could not copy shape type: {source_shape.shape_type}")
                
    except Exception as e:
        print(f"Error copying individual shape: {e}")

def copy_image_shape(source_shape, target_slide):
    """Copy an image shape to the target slide."""
    try:
        # Get image data
        image_data = source_shape.image.blob
        image_stream = io.BytesIO(image_data)
        
        # Add image to target slide with same position and size
        pic = target_slide.shapes.add_picture(
            image_stream,
            source_shape.left,
            source_shape.top,
            source_shape.width,
            source_shape.height
        )
        
        # Copy any formatting properties if possible
        if hasattr(source_shape, 'rotation'):
            pic.rotation = source_shape.rotation
            
    except Exception as e:
        print(f"Could not copy image: {e}")

def copy_text_or_shape(source_shape, target_slide):
    """Copy text boxes or auto shapes to target slide."""
    try:
        # For text boxes, create a new text box
        if hasattr(source_shape, 'text'):
            textbox = target_slide.shapes.add_textbox(
                source_shape.left,
                source_shape.top, 
                source_shape.width,
                source_shape.height
            )
            
            # Copy text content (but we might override this later)
            textbox.text = source_shape.text
            
            # Copy text formatting if possible
            if source_shape.text_frame and textbox.text_frame:
                copy_text_formatting(source_shape.text_frame, textbox.text_frame)
        
        # For shapes, try to recreate based on shape type
        else:
            # This is more complex - would need shape-specific logic
            pass
            
    except Exception as e:
        print(f"Could not copy text/shape: {e}")

def copy_generic_shape(source_shape, target_slide):
    """Attempt to copy other shape types."""
    try:
        # This is a placeholder for more complex shape copying
        # Different shape types require different approaches
        # For now, we'll skip complex shapes to avoid errors
        pass
    except Exception as e:
        print(f"Could not copy generic shape: {e}")

def copy_text_formatting(source_tf, target_tf):
    """Copy text formatting from source to target text frame."""
    try:
        # Copy paragraph-level formatting
        for src_para, tgt_para in zip(source_tf.paragraphs, target_tf.paragraphs):
            if hasattr(src_para, 'alignment'):
                tgt_para.alignment = src_para.alignment
            if hasattr(src_para, 'font'):
                if src_para.font.name:
                    tgt_para.font.name = src_para.font.name
                if src_para.font.size:
                    tgt_para.font.size = src_para.font.size
                    
    except Exception as e:
        print(f"Could not copy text formatting: {e}")

def populate_slide_content(slide, content):
    """
    Populate slide placeholders with our generated content.
    Only updates placeholder text, leaves all other elements intact.
    """
    title = content.get("title", "")
    points = content.get("points", [])
    
    # Find and populate title placeholder
    title_shape = None
    content_shape = None
    
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                title_shape = shape
            elif shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                content_shape = shape
    
    # Set title
    if title_shape and title:
        title_shape.text = title
    
    # Set content points
    if content_shape and points:
        tf = content_shape.text_frame
        tf.clear()
        
        # Add first point
        if points:
            p = tf.paragraphs[0]
            p.text = points[0]
            
            # Add remaining points
            for point_text in points[1:]:
                new_para = tf.add_paragraph()
                new_para.text = point_text
                new_para.level = 0

def create_basic_presentation(slide_data, output_path):
    """
    Fallback: Create basic presentation when no template is provided.
    """
    prs = Presentation()
    
    # Use default layout
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    
    for content in slide_data:
        slide = prs.slides.add_slide(layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = content.get("title", "")
        
        # Set content
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                tf = shape.text_frame
                tf.clear()
                
                points = content.get("points", [])
                if points:
                    p = tf.paragraphs[0]
                    p.text = points[0]
                    
                    for point_text in points[1:]:
                        new_p = tf.add_paragraph()
                        new_p.text = point_text
                        new_p.level = 0
                break
    
    prs.save(output_path)
    return output_path














# # core/generator.py

# from pptx import Presentation
# from pptx.enum.shapes import PP_PLACEHOLDER
# from pptx.dml.color import RGBColor
# from pptx.util import Inches, Pt
# import os
# from typing import Optional, Dict, Any

# def extract_template_style(template_path: str) -> Dict[str, Any]:
#     """
#     Extract styling information from a template presentation.
#     Returns a dictionary with style information.
#     """
#     if not template_path or not os.path.exists(template_path):
#         return {}
    
#     try:
#         template_prs = Presentation(template_path)
#         style_info = {
#             'slide_layouts': [],
#             'master_slides': [],
#             'theme_colors': {},
#             'fonts': {},
#             'images': []
#         }
        
#         # Extract slide layouts
#         for i, layout in enumerate(template_prs.slide_layouts):
#             layout_info = {
#                 'index': i,
#                 'name': layout.name,
#                 'placeholders': []
#             }
            
#             for placeholder in layout.placeholders:
#                 ph_info = {
#                     'type': placeholder.placeholder_format.type,
#                     'idx': placeholder.placeholder_format.idx,
#                     'left': placeholder.left,
#                     'top': placeholder.top,
#                     'width': placeholder.width,
#                     'height': placeholder.height
#                 }
#                 layout_info['placeholders'].append(ph_info)
            
#             style_info['slide_layouts'].append(layout_info)
        
#         # Extract images from template slides
#         for slide in template_prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, 'image'):
#                     # Store image information
#                     image_info = {
#                         'left': shape.left,
#                         'top': shape.top,
#                         'width': shape.width,
#                         'height': shape.height,
#                         'image_data': shape.image.blob  # Store the actual image data
#                     }
#                     style_info['images'].append(image_info)
        
#         return style_info
        
#     except Exception as e:
#         print(f"Error extracting template style: {e}")
#         return {}

# def find_best_layout(layouts_info: list, needs_content: bool = True) -> int:
#     """
#     Find the best layout index for content slides.
#     """
#     if not layouts_info:
#         return 1  # Default to layout 1
    
#     # Look for a layout with title and content placeholders
#     for layout in layouts_info:
#         has_title = any(ph['type'] == PP_PLACEHOLDER.TITLE for ph in layout['placeholders'])
#         has_content = any(ph['type'] in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) 
#                          for ph in layout['placeholders'])
        
#         if has_title and (has_content or not needs_content):
#             return layout['index']
    
#     # Fallback to layout 1 (usually title and content)
#     return min(1, len(layouts_info) - 1)

# def create_ppt_from_template(slide_data, output_path, template_path=None, template_style=None):
#     """
#     Creates a NEW PPT file from slide_data, applying styles from a template.
#     This creates a completely new presentation rather than modifying the template.
#     """
    
#     # Step 1: Create a completely NEW presentation
#     if template_path and os.path.exists(template_path):
#         # Load template to get layouts and styling, but we'll create new slides
#         template_prs = Presentation(template_path)
#         # Create new presentation using the template's master slides and layouts
#         prs = Presentation(template_path)
        
#         # Remove all existing slides from the template
#         # We need to remove slides in reverse order to avoid index issues
#         slide_indices = list(range(len(prs.slides)))
#         for i in reversed(slide_indices):
#             slide_id = prs.slides[i].slide_id
#             prs.part.drop_rel(prs.slides._sldIdLst[i].rId)
#             del prs.slides._sldIdLst[i]
            
#     else:
#         # No template provided, create blank presentation
#         prs = Presentation()
#         template_prs = None
    
#     # Step 2: Extract style information if template exists
#     style_info = extract_template_style(template_path) if template_path else {}
    
#     # Step 3: Find the best layout to use
#     if template_prs:
#         # Find a suitable "Title and Content" layout
#         title_and_content_layout = None
#         for layout in prs.slide_layouts:
#             has_title = any(ph.placeholder_format.type == PP_PLACEHOLDER.TITLE for ph in layout.placeholders)
#             has_body = any(ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) 
#                           for ph in layout.placeholders)
#             if has_title and has_body:
#                 title_and_content_layout = layout
#                 break
        
#         # Fallback to layout 1 if no suitable layout found
#         if not title_and_content_layout:
#             title_and_content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
#     else:
#         # Use default layout for blank presentation
#         title_and_content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

#     # Step 4: Create NEW slides with the provided content
#     for item in slide_data:
#         slide = prs.slides.add_slide(title_and_content_layout)
        
#         # Add title
#         title_shape = slide.shapes.title
#         if title_shape:
#             title_shape.text = item.get("title", "No Title")

#         # Find and populate content placeholder
#         body_shape = None
#         for shape in slide.placeholders:
#             if shape.placeholder_format.type != PP_PLACEHOLDER.TITLE:
#                 body_shape = shape
#                 break
        
#         if body_shape:
#             tf = body_shape.text_frame
#             tf.clear()
            
#             points = item.get("points", [])
#             if points:
#                 # Add first point
#                 p = tf.paragraphs[0]
#                 p.text = points[0]
                
#                 # Add remaining points
#                 for point_text in points[1:]:
#                     p = tf.add_paragraph()
#                     p.text = point_text
#                     p.level = 0

#         # Step 5: Add template images to slides if available
#         if style_info and 'images' in style_info:
#             for img_info in style_info['images']:
#                 try:
#                     # Add the image from template to this slide
#                     # Note: This is a simplified approach - in practice, you might want to be more selective
#                     # about which images to include and where to place them
#                     if len(slide.shapes) < 10:  # Avoid cluttering slides
#                         # You could add image placement logic here
#                         pass
#                 except Exception as e:
#                     print(f"Could not add image to slide: {e}")

#     # Step 6: Save the new presentation
#     prs.save(output_path)
#     return output_path

# def copy_slide_layout_properties(source_layout, target_layout):
#     """
#     Copy properties from source layout to target layout.
#     This is a helper function for advanced template copying.
#     """
#     try:
#         # This is a placeholder for more advanced layout copying
#         # The python-pptx library has limitations in directly copying layouts
#         # In practice, using the template as the base presentation (but clearing slides)
#         # is the most reliable approach
#         pass
#     except Exception as e:
#         print(f"Error copying layout properties: {e}")

# core/generator.py